# utils.py
# Yordamchi funksiyalar: faylni o'qish va prompt yaratish

import io
import sys 

# Mavjud kutubxonalar
try:
    import docx
except ImportError:
    print("XATOLIK: 'python-docx' kutubxonasi o'rnatilmagan. (pip install python-docx)")
    sys.exit(1)

# --- YANGI KUTUBXONALAR ---
try:
    import fitz  # PyMuPDF
except ImportError:
    print("XATOLIK: 'PyMuPDF' kutubxonasi o'rnatilmagan. (pip install PyMuPDF)")
    sys.exit(1)

try:
    from pptx import Presentation
except ImportError:
    print("XATOLIK: 'python-pptx' kutubxonasi o'rnatilmagan. (pip install python-pptx)")
    sys.exit(1)
# -------------------------


# --- get_text_from_file FUNKSIYASI (PDF va PPTX qo'shilgan) ---
def get_text_from_file(file_storage):
    """
    Flask'ning FileStorage obyektini qabul qilib, undan matnni qaytaradi.
    .txt, .docx, .pdf, va .pptx fayllarni qo'llab-quvvatlaydi.
    """
    filename = file_storage.filename
    
    if filename.endswith('.docx'):
        try:
            doc = docx.Document(io.BytesIO(file_storage.read()))
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return '\n'.join(full_text)
        except Exception as e:
            print(f".docx faylni o'qishda xatolik: {e}")
            raise ValueError(f".docx faylni o'qishda xatolik: {e}")
            
    elif filename.endswith('.txt'):
        try:
            text = file_storage.read().decode('utf-8')
            return text
        except Exception as e:
            print(f".txt faylni o'qishda xatolik: {e}")
            raise ValueError(f".txt faylni o'qishda xatolik: {e}")

    elif filename.endswith('.pdf'):
        try:
            doc = fitz.open(stream=file_storage.read(), filetype="pdf")
            full_text = []
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                full_text.append(page.get_text())
            doc.close()
            return '\n'.join(full_text)
        except Exception as e:
            print(f".pdf faylni o'qishda xatolik: {e}")
            raise ValueError(f".pdf faylni o'qishda xatolik: {e}")

    elif filename.endswith('.pptx'):
        try:
            ppt = Presentation(io.BytesIO(file_storage.read()))
            full_text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            return '\n'.join(full_text)
        except Exception as e:
            print(f".pptx faylni o'qishda xatolik: {e}")
            raise ValueError(f".pptx faylni o'qishda xatolik: {e}")
            
    else:
        raise ValueError(f"Fayl formati qo'llab-quvvatlanmaydi: {filename}. Faqat .txt, .docx, .pdf, .pptx.")

# --- create_dynamic_prompt FUNKSIYASI (RAQAMLASH TUZATILDI) ---
def create_dynamic_prompt(q_type, q_difficulty, q_count, test_format, output_language, lecture_text):
    """
    Foydalanuvchi sozlamalari va matn asosida model uchun dinamik prompt yaratadi.
    """
    
    # Tilga bog'liq matnlar
    if output_language == 'ru':
        lang_instructions = {
            "role": "Вы - ИИ-помощник для преподавателей университетов Узбекистана.",
            "task": "Ваша задача - проанализировать предоставленный текст лекции и сгенерировать вопросы для проверки знаний студентов в соответствии с указанными требованиями.",
            "source_text_header": "ТЕКСТ:",
            "requirements_header": "ТРЕБОВАНИЯ, которые необходимо выполнить:",
            "q_type_label": "1. ТИП ВОПРОСОВ:",
            "difficulty_label": "2. СЛОЖНОСТЬ:",
            "count_label": "3. КОЛИЧЕСТВО ВОПРОСОВ:",
            "test_format_label": "4. ФОРМАТ ТЕСТОВ:",
            "based_on_text": "Все вопросы должны основываться только и исключительно на содержании предоставленного выше текста.",
            # --- MANA SHU QISM YUMSHATILDI (RAQAMLARGA RUXSAT BERILDI) ---
            "no_reference": """ПРИ СОСТАВЛЕНИИ ВОПРОСОВ КАТЕГОРИЧЕСКИ ЗАПРЕЩАЕТСЯ:
1. ЛЮБЫЕ фразы, ссылающиеся на источник, текст, таблицу, рисунок, пример, приложение или раздел.
Например: "согласно источнику", "на основе текста", "по таблице", "согласно примеру 1.3", "на рисунке выше", "согласно сравнительной таблице" и т.п.
2. НЕ начинайте текст вопроса со слов "Вопрос:", "Текст вопроса:", "Тест:", "Тестовый вопрос:".
Вопросы ДОЛЖНЫ быть полностью самостоятельными, и для ответа на них должен быть достаточен только сам текст вопроса.""",
            # -----------------------------------------------------------
            "output_format_header": "Верните ответ в следующем формате (включайте только запрошенные типы):",
            "final_instruction": f"ВЕСЬ ОТВЕТ ДОЛЖЕН БЫТЬ СГЕНЕРИРОВАН НА РУССКОМ ЯЗЫКЕ.",
            # --- RAQAMLASH BUYRUG'I QO'SHILDI ---
            "type_all": "Сгенерируйте смесь теоретических (открытых) и тестовых (закрытых, с вариантами) вопросов. Теоретические вопросы пронумеруйте (1., 2., 3...).",
            "type_theory": "Сгенерируйте только теоретические (открытые) вопросы. Пронумеруйте их (1., 2., 3...). НЕ СОЗДАВАЙТЕ ТЕСТОВЫЕ ВОПРОСЫ.",
            # --------------------------------------
            "type_test": "Сгенерируйте только тестовые (закрытые, с вариантами) вопросы. НЕ СОЗДАВАЙТЕ ТЕОРЕТИЧЕСКИЕ ВОПРОСЫ.",
            "difficulty_easy": "Вопросы должны быть легкими, сосредоточенными только на основных терминах и определениях.",
            "difficulty_medium": "Вопросы должны быть средней сложности.",
            "difficulty_hard": "Вопросы должны быть сложными, требующими анализа и сравнения.",
            "count_text": f"Сгенерируйте примерно {q_count} вопросов(а) всего.",
            "format_theory_header": "### I. Теоретические (Открытые) Вопросы",
            "format_test_header": "### II. Тестовые (Закрытые) Вопросы",
            "hemis_instruction": """
    ФОРМАТ ДЛЯ ТЕСТОВЫХ ВОПРОСОВ: Для тестовых вопросов используйте *только* и *строго* следующий специальный формат HEMIS:
    Текст вопроса 1 (БЕЗ каких-либо надписей "Вопрос:" или нумерации "1.")
    ====
    Вариант А
    ====
    #Правильный вариант ответа
    ====
    Вариант C
    ====
    Вариант D
    ++++
    Текст вопроса 2
    ... и так далее.
    ПРАВИЛА (ОЧЕНЬ ВАЖНО!):
    1. Текст вопроса, каждый вариант и '====' должны быть на отдельной строке.
    2. Перед правильным вариантом ответа ставьте *только* знак '#'.
    3. Разделяйте блоки вопросов *только* знаком '++++'.
    4. НЕ ИСПОЛЬЗУЙТЕ никакой другой нумерации (1., 2.) или буквенных обозначений (а), б)).
    5. НЕ ИСПОЛЬЗУЙТЕ скобки ([ ]) вокруг вариантов или текста вопроса.
    6. НЕ начинайте текст вопроса со слов "Вопрос:", "Текст вопроса:", "Тест:" или с цифры. Сразу пишите сам текст вопроса.
    """,
            # --- RAQAMLASH BUYRUG'I QO'SHILDI ---
            "normal_test_format": "ФОРМАТ ДЛЯ ТЕСТОВЫХ ВОПРОСОВ: Создайте тестовые вопросы со стандартной нумерацией (1., 2., 3...) и вариантами (а, б, в, г)."
            # --------------------------------------
        }
    else: # output_language == 'uz' (Standart)
        lang_instructions = {
            "role": "Siz O'zbekistondagi universitet o'qituvchilari uchun yordamchi AI'siz.",
            "task": "Vazifangiz - taqdim etilgan ma'ruza matnini tahlil qilib, ko'rsatilgan talablar asosida savollar yaratish.",
            "source_text_header": "MATN:",
            "requirements_header": "bajarilishi shart bo'lgan KO'RSATMALAR:",
            "q_type_label": "1. SAVOL TURI:",
            "difficulty_label": "2. MURAKKABLIK:",
            "count_label": "3. SAVOLLAR SONI:",
            "test_format_label": "4. TEST FORMATI:",
            "based_on_text": "Barcha savollar faqat va faqat yuqorida berilgan matn mazmuniga asoslansin.",
            # --- MANA SHU QISM YUMSHATILDI (RAQAMLARGA RUXSAT BERILDI) ---
            "no_reference": """SAVOLLARNI TUZISHDA QUYIDAGILAR QAT'IYAN MAN ETILADI:
1. Manbaga, matnga, jadvalga, rasmga, misolga, ilovaga yoki bo'limga ishora qiluvchi, havola beruvchi HAR QANDAY ibora.
Masalan: "manbaga ko'ra", "matnga asoslanib", "jadvalga ko'ra", "1.3-misolga asosan", "yuqoridagi rasmda", "taqqoslash jadvaliga ko'ra" va hokazo.
2. Savol matnini "Savol:", "Savol matni:", "Test:", "Test savoli:" kabi so'zlar bilan BOSHLAMANG (HEMIS formati bundan mustasno emas).
Savollar to'liq mustaqil bo'lishi va ularga javob berish uchun faqat savol matnining o'zi yetarli bo'lishi SHART.""",
            # -----------------------------------------------------------
            "output_format_header": "Javobni quyidagi formatda qaytaring (faqat so'ralgan turlarni kiriting):",
            "final_instruction": "", 
            # --- RAQAMLASH BUYRUG'I QO'SHILDI ---
            "type_all": "Nazariy (ochiq) va Test (yopiq, variantli) savollarni aralash yarating. Nazariy savollarni raqamlang (1., 2., 3...).",
            "type_theory": "Faqat Nazariy (ochiq) savollar yarating. Ularni raqamlang (1., 2., 3...). TEST SAVOLLARI YARATMA.",
            # --------------------------------------
            "type_test": "Faqat Test (yopiq, variantli) savollar yarating. NAZARIY SAVOLLAR YARATMA.",
            "difficulty_easy": "Savollar oson, faqat asosiy terminlar va ta'riflarga e'tibor qaratsin.",
            "difficulty_medium": "Savollar o'rta murakkablikda bo'lsin.",
            "difficulty_hard": "Savollar qiin, tahlil qilishni va taqqoslashni talab qiladigan bo'lsin.",
            "count_text": f"Jami taxminan {q_count} ta savol generatsiya qiling.",
            "format_theory_header": "### I. Nazariy (Ochiq) Savollar",
            "format_test_header": "### II. Test (Yopiq) Savollari",
            "hemis_instruction": """
    TEST SAVOLLARI UCHUN FORMAT: Test savollari uchun *faqat* va *aniq* quyidagi maxsus HEMIS formatidan foydalan:
    Savol matni 1 (HECH QANDAY "Savol:" degan yozuvlarsiz yoki "1." raqamisiz)
    ====
    Variant A
    ====
    #To'g'ri javob varianti
    ====
    Variant C
    ====
    Variant D
    ++++
    Savol matni 2
    ... va hokazo.
    QOIDALAR (JUDA MUHIM!):
    1. Savol matni, har bir variant va '====' bo'luvchisi alohida qatorda bo'lsin.
    2. To'g'ri javob variantidan oldin faqat '#' belgisi ishlatilsin.
    3. Har bir savol bloki orasini faqat '++++' belgisi bilan ajrat.
    4. Boshqa hech qanday raqamlash (1., 2.) yoki (a), b)) kabi harf bilan belgilash ISHLATILMASIN.
    5. Variantlar yoki savol matni atrofida boshqa belgilar ISHLATILMASIN.
    6. Savol matnini "Savol:", "Savol matni:", "Test:" kabi so'zlar bilan yoki raqam bilan BOSHLAMANG. To'g'ridan-to'g'ri savolning o'zini yozing.
    """,
             # --- RAQAMLASH BUYRUG'I QO'SHILDI ---
            "normal_test_format": "TEST SAVOLLARI UCHUN FORMAT: Test savollarini standart raqamlash (1., 2., 3...) va variantlar (a, b, c, d) bilan yarating."
            # --------------------------------------
        }

    # Prompt qismlarini shakllantirish
    type_instruction = ""
    format_instruction = ""
    test_format_instruction = ""

    if q_type != 'faqat_nazariy':
        if test_format == 'hemis':
            test_format_instruction = lang_instructions["hemis_instruction"]
        else:
            test_format_instruction = lang_instructions["normal_test_format"]

    if q_type == 'faqat_nazariy':
        type_instruction = lang_instructions["type_theory"]
        format_instruction = lang_instructions["format_theory_header"]
    elif q_type == 'faqat_test':
        type_instruction = lang_instructions["type_test"]
        if test_format == 'hemis':
             format_instruction = ""
        else:
            format_instruction = lang_instructions["format_test_header"]
    else: # 'hammasi'
        type_instruction = lang_instructions["type_all"]
        if test_format == 'hemis':
             format_instruction = f"{lang_instructions['format_theory_header']}\n\n"
        else:
            format_instruction = f"{lang_instructions['format_theory_header']}\n\n{lang_instructions['format_test_header']}"

    # Kalit xatoligini oldini olish uchun tuzatilgan mantiq
    difficulty_key = f"difficulty_{q_difficulty}" 
    if difficulty_key not in lang_instructions:
        difficulty_key = "difficulty_medium"
    
    difficulty_instruction = lang_instructions[difficulty_key]
    count_instruction = lang_instructions["count_text"]

    # --- YAKUNIY PROMPT (TILGA MOSLASHTIRILGAN) ---
    prompt = f"""{lang_instructions["role"]}
{lang_instructions["task"]}

{lang_instructions["source_text_header"]}
---
{lecture_text}
---

{lang_instructions["requirements_header"]}
{lang_instructions["q_type_label"]} {type_instruction}
{lang_instructions["difficulty_label"]} {difficulty_instruction}
{lang_instructions["count_label"]} {count_instruction}
{lang_instructions["test_format_label"]} {test_format_instruction}

{lang_instructions["based_on_text"]}
{lang_instructions["no_reference"]}

{lang_instructions["output_format_header"]}
{format_instruction}

{lang_instructions["final_instruction"]}
    """
    
    return prompt
